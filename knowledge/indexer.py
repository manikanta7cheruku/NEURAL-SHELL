"""
=============================================================================
PROJECT SEVEN - knowledge/indexer.py (Document Indexer)
Version: 2.0

Supports: .txt, .md, .pdf, .docx, .pptx, .xlsx
=============================================================================
"""

import os
import json
import hashlib
from colorama import Fore
from knowledge.core import store_chunk
import config

KNOWLEDGE_DIR  = os.path.join("seven_data", "knowledge")
CUSTOM_DIR     = os.path.join(KNOWLEDGE_DIR, "custom")
MANIFEST_FILE  = os.path.join(KNOWLEDGE_DIR, "index_manifest.json")

os.makedirs(CUSTOM_DIR, exist_ok=True)

SUPPORTED_EXTENSIONS = {".txt", ".md", ".pdf", ".docx", ".pptx", ".xlsx"}


# ── Manifest ──────────────────────────────────────────────────────────────

def _load_manifest():
    if os.path.exists(MANIFEST_FILE):
        try:
            with open(MANIFEST_FILE, 'r') as f:
                return json.load(f)
        except Exception:
            return {}
    return {}


def _save_manifest(manifest):
    try:
        with open(MANIFEST_FILE, 'w') as f:
            json.dump(manifest, f, indent=2)
    except Exception as e:
        print(Fore.RED + f"[INDEXER] Manifest save error: {e}")


def _file_hash(filepath):
    h = hashlib.md5()
    try:
        with open(filepath, 'rb') as f:
            for chunk in iter(lambda: f.read(8192), b''):
                h.update(chunk)
        return h.hexdigest()
    except Exception:
        return None


def get_index_manifest():
    return _load_manifest()


# ── Text Extraction ────────────────────────────────────────────────────────

def _extract_txt(filepath):
    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()


def _extract_pdf(filepath):
    try:
        import PyPDF2
        text = []
        with open(filepath, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                try:
                    t = page.extract_text()
                    if t:
                        text.append(t)
                except Exception:
                    pass
        return "\n\n".join(text)
    except ImportError:
        raise ImportError("PyPDF2 not installed. Run: pip install pypdf2")


def _extract_docx(filepath):
    try:
        import docx
        doc = docx.Document(filepath)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    paragraphs.append(row_text)
        return "\n\n".join(paragraphs)
    except ImportError:
        raise ImportError("python-docx not installed. Run: pip install python-docx")


def _extract_pptx(filepath):
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                slides.append(f"[Slide {i}]\n" + "\n".join(slide_text))
        return "\n\n".join(slides)
    except ImportError:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")


def _extract_xlsx(filepath):
    try:
        import openpyxl
        wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        sheets = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None and str(c).strip()]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(sheets)
    except ImportError:
        raise ImportError("openpyxl not installed. Run: pip install openpyxl")


def extract_text(filepath):
    """
    Extract text from any supported file type.
    Returns raw text string.
    """
    ext = os.path.splitext(filepath)[1].lower()
    extractors = {
        ".txt":  _extract_txt,
        ".md":   _extract_txt,
        ".pdf":  _extract_pdf,
        ".docx": _extract_docx,
        ".pptx": _extract_pptx,
        ".xlsx": _extract_xlsx,
    }
    if ext not in extractors:
        raise ValueError(f"Unsupported format: {ext}")
    return extractors[ext](filepath)


# ── Chunking ───────────────────────────────────────────────────────────────

def _chunk_text(text, chunk_size=None):
    if chunk_size is None:
        chunk_size = config.KEY.get("knowledge", {}).get("chunk_size", 600)

    paragraphs = text.split("\n\n")
    chunks = []
    current = ""

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue
        if len(current) + len(para) + 2 <= chunk_size:
            current += para + "\n\n"
        else:
            if current.strip():
                chunks.append(current.strip())
            if len(para) > chunk_size:
                sentences = para.replace(". ", ".\n").split("\n")
                current = ""
                for sent in sentences:
                    sent = sent.strip()
                    if not sent:
                        continue
                    if len(current) + len(sent) + 1 <= chunk_size:
                        current += sent + " "
                    else:
                        if current.strip():
                            chunks.append(current.strip())
                        current = sent + " "
            else:
                current = para + "\n\n"

    if current.strip():
        chunks.append(current.strip())

    return chunks


# ── Indexing ───────────────────────────────────────────────────────────────

def index_file(filepath):
    """
    Index a single file. Returns (success, chunks_added, message).
    """
    filepath = os.path.abspath(filepath)

    if not os.path.exists(filepath):
        return False, 0, f"File not found: {filepath}"

    ext = os.path.splitext(filepath)[1].lower()
    if ext not in SUPPORTED_EXTENSIONS:
        return False, 0, f"Unsupported format: {ext}. Supported: {', '.join(SUPPORTED_EXTENSIONS)}"

    manifest  = _load_manifest()
    file_hash = _file_hash(filepath)
    filename  = os.path.basename(filepath)

    if filename in manifest and manifest[filename].get("hash") == file_hash:
        return True, 0, f"Already indexed: {filename} (unchanged)"

    # Extract text
    try:
        text = extract_text(filepath)
    except ImportError as e:
        return False, 0, str(e)
    except Exception as e:
        return False, 0, f"Read error: {e}"

    if not text.strip():
        return False, 0, f"No text content found in: {filename}"

    # Chunk and store
    chunks = _chunk_text(text)
    if not chunks:
        return False, 0, f"No content to index in {filename}"

    # Remove old chunks if re-indexing changed file
    if filename in manifest:
        old_hash = manifest[filename].get("hash", "")
        if old_hash and old_hash != file_hash:
            old_chunks = manifest[filename].get("chunks", 0)
            for i in range(old_chunks):
                try:
                    from knowledge.core import knowledge_collection
                    knowledge_collection.delete(ids=[f"{filename}_chunk_{i}_{old_hash[:8]}"])
                except Exception:
                    pass

    stored = 0
    for i, chunk in enumerate(chunks):
        chunk_id = f"{filename}_chunk_{i}_{file_hash[:8]}"
        store_chunk(
            text=chunk,
            chunk_id=chunk_id,
            source=filename,
            category="document"
        )
        stored += 1

    manifest[filename] = {
        "hash":   file_hash,
        "chunks": stored,
        "path":   filepath,
        "ext":    ext,
        "size_kb": round(os.path.getsize(filepath) / 1024, 1),
    }
    _save_manifest(manifest)

    print(Fore.GREEN + f"[INDEXER] Indexed {filename}: {stored} chunks")
    return True, stored, f"Indexed {filename}: {stored} chunks"


def remove_file(filename):
    """
    Remove a file from the knowledge base.
    Returns (success, message).
    """
    manifest = _load_manifest()
    if filename not in manifest:
        return False, f"File not in index: {filename}"

    entry     = manifest[filename]
    file_hash = entry.get("hash", "")
    chunks    = entry.get("chunks", 0)

    removed = 0
    for i in range(chunks):
        try:
            from knowledge.core import knowledge_collection
            knowledge_collection.delete(ids=[f"{filename}_chunk_{i}_{file_hash[:8]}"])
            removed += 1
        except Exception:
            pass

    del manifest[filename]
    _save_manifest(manifest)

    # Also delete physical file from custom dir if it exists
    physical = os.path.join(CUSTOM_DIR, filename)
    if os.path.exists(physical):
        try:
            os.remove(physical)
        except Exception:
            pass

    print(Fore.GREEN + f"[INDEXER] Removed {filename}: {removed} chunks deleted")
    return True, f"Removed {filename}"


def index_directory(dirpath=None):
    if dirpath is None:
        dirpath = CUSTOM_DIR

    if not os.path.exists(dirpath):
        os.makedirs(dirpath, exist_ok=True)
        return 0, 0, ["Directory created. Drop files there."]

    total_files  = 0
    total_chunks = 0
    messages     = []

    for filename in os.listdir(dirpath):
        ext = os.path.splitext(filename)[1].lower()
        if ext not in SUPPORTED_EXTENSIONS:
            continue
        filepath = os.path.join(dirpath, filename)
        success, chunks, msg = index_file(filepath)
        messages.append(msg)
        if success and chunks > 0:
            total_files  += 1
            total_chunks += chunks

    if not messages:
        messages.append(f"No supported files found in {dirpath}.")

    return total_files, total_chunks, messages